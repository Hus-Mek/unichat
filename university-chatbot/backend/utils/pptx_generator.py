"""PowerPoint presentation generator for university deployment proposal."""

import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .cost_calculator import (
    aws_cloud_costs,
    azure_cloud_costs,
    comparison_table,
    groq_hybrid_costs,
    on_premise_costs,
)

# Colour palette
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _add_title_slide(prs: Presentation, uni_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI-Powered University Chatbot"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"Deployment Proposal for {uni_name}"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "RAG-Based Knowledge Assistant with Multi-Level Access Control"
    p3.font.size = Pt(16)
    p3.font.color.rgb = DARK_TEXT
    p3.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(12)


def _add_cost_chart(prs: Presentation, daily_queries: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "5-Year Total Cost of Ownership Comparison"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    table_data = comparison_table(daily_queries)
    chart_data = CategoryChartData()
    chart_data.categories = [row["option"] for row in table_data]
    chart_data.add_series("Year 1", [row["year_1"] for row in table_data])
    chart_data.add_series("Year 2+/yr", [row["year_2_plus"] for row in table_data])
    chart_data.add_series("5-Year TCO", [row["five_year_tco"] for row in table_data])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False


def _add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5 + num_rows * 0.6),
    )
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    for r, row_data in enumerate(rows, start=1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GREY


def generate_proposal(
    university_name: str = "University",
    expected_daily_queries: int = 1000,
    preferred_deployment: str | None = None,
) -> io.BytesIO:
    """Generate the full proposal PowerPoint and return as BytesIO."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    _add_title_slide(prs, university_name)

    # Slide 2: Executive Summary
    _add_content_slide(prs, "Executive Summary", [
        "AI-powered chatbot trained exclusively on university documents",
        "5-level access control: Public, Student, Faculty, Admin Staff, Executive Board",
        "Multi-tenant architecture supporting per-department document collections",
        "RAG (Retrieval-Augmented Generation) ensures answers come only from approved documents",
        "Flexible deployment: on-premise, cloud, or hybrid API-based",
        f"Designed for ~{expected_daily_queries:,} queries per day",
    ])

    # Slide 3: Current Demo vs Production
    _add_content_slide(prs, "Current Demo vs. Production System", [
        "DEMO: No real authentication - access level selected via dropdown",
        "DEMO: Only 3 access levels (public, student, faculty)",
        "DEMO: Single document collection, no department separation",
        "DEMO: No audit logging, session history, or usage tracking",
        "PRODUCTION: JWT authentication with SSO-ready architecture",
        "PRODUCTION: 5 access levels with document and collection-level control",
        "PRODUCTION: Multi-tenant with per-department document spaces",
        "PRODUCTION: Full audit trail, session persistence, cost tracking",
    ])

    # Slide 4: Architecture
    _add_content_slide(prs, "System Architecture", [
        "Frontend: Streamlit web interface (upgradeable to React)",
        "Backend: FastAPI REST API with async request handling",
        "Database: PostgreSQL for users, sessions, audit logs",
        "Vector Store: ChromaDB for semantic document search",
        "Embeddings: all-MiniLM-L6-v2 (384 dimensions)",
        "LLM: Pluggable provider (Groq API / vLLM / Ollama / Azure OpenAI)",
        "Deployment: Docker Compose with nginx reverse proxy",
    ])

    # Slide 5: Key Features
    _add_content_slide(prs, "Key Features", [
        "Context-only answers: strict guardrails prevent hallucination",
        "Document upload: PDF, Excel, Word with automatic text extraction",
        "Semantic search: finds relevant passages even with paraphrased queries",
        "Cost tracking: per-query token usage and cost monitoring",
        "Session history: persistent chat sessions for each user",
        "Admin dashboard: user management, audit logs, system statistics",
        "Report generation: automated PowerPoint proposals with cost analysis",
    ])

    # Slide 6: Access Control
    _add_table_slide(prs, "5-Level Access Control",
        ["Level", "Can Access", "Typical Users"],
        [
            ["Public", "Public documents only", "Anonymous / visitors"],
            ["Student", "Public + Student docs", "Enrolled students"],
            ["Faculty", "Public + Student + Faculty docs", "Professors, lecturers"],
            ["Admin Staff", "All except Executive", "Department admins, IT staff"],
            ["Executive Board", "All documents", "University leadership"],
        ],
    )

    # Slide 7: Multi-Tenancy
    _add_content_slide(prs, "Multi-Tenant Document Management", [
        "Each department can have its own document collection",
        "Collections have a minimum access level requirement",
        "Fine-grained grants: specific users can access specific collections",
        "Documents within collections have their own access levels",
        "Two-layer security: collection access + document access must both pass",
        "Example: HR collection (admin_staff) with salary docs (executive_board)",
    ])

    # Slide 8: RAG Pipeline
    _add_content_slide(prs, "RAG Pipeline: How Answers Are Generated", [
        "1. User submits a question through the chat interface",
        "2. Question is converted to a 384-dimensional embedding vector",
        "3. ChromaDB performs semantic similarity search across accessible documents",
        "4. Access control filters ensure user only sees permitted content",
        "5. Top 15 relevant passages are assembled into context",
        "6. LLM generates an answer using ONLY the provided context",
        "7. Response includes source citations, token usage, and cost",
    ])

    # Slide 9: Security
    _add_content_slide(prs, "Security & SSO Readiness", [
        "JWT-based authentication with configurable token expiry",
        "bcrypt password hashing for local accounts",
        "SSO integration points: SAML 2.0, OpenID Connect, LDAP",
        "All actions logged in audit trail (login, query, upload, delete)",
        "Rate limiting per user to prevent API abuse",
        "Document-level access control enforced at query time",
        "System prompt prevents LLM from leaking confidential information",
    ])

    # Slide 10: Cost Comparison Chart
    _add_cost_chart(prs, expected_daily_queries)

    # Slide 11: On-Prem Details
    onprem = on_premise_costs("mid")
    _add_table_slide(prs, "Option A: On-Premise GPU Server",
        ["Component", "Specification", "Cost"],
        [
            ["GPU", "2x NVIDIA A100 40GB", "$30,000"],
            ["Server (CPU, RAM, Storage)", "AMD EPYC + 256GB + 4TB NVMe", "$20,000"],
            ["Annual Operations", "Electricity + IT maintenance", "$10,000/yr"],
            ["Software", "vLLM + open-source models", "$0"],
            ["Year 1 Total", "", f"${onprem['yearly'].year_1:,.0f}"],
            ["5-Year TCO", "", f"${onprem['yearly'].five_year_tco:,.0f}"],
        ],
    )

    # Slide 12: Cloud Details
    azure = azure_cloud_costs(expected_daily_queries)
    _add_table_slide(prs, "Option B: Azure Cloud Deployment",
        ["Resource", "Specification", "Monthly Cost"],
        [
            ["GPU VM", "NC24ads A100 v4", f"${azure['breakdown']['GPU VM (NC24ads A100)']:,.0f}"],
            ["App VM", "D4s v5 (4 vCPU, 16GB)", f"${azure['breakdown']['App VM (D4s v5)']:,.0f}"],
            ["PostgreSQL", "General Purpose, 4 vCores", f"${azure['breakdown']['PostgreSQL (4 vCores)']:,.0f}"],
            ["Storage + Network", "100GB Blob + VNet", f"${azure['breakdown']['Blob Storage'] + azure['breakdown']['Networking']:,.0f}"],
            ["Monthly Total", "", f"${azure['monthly_cost']:,.0f}"],
            ["Annual Total", "", f"${azure['yearly'].year_1:,.0f}"],
        ],
    )

    # Slide 13: Hybrid Details
    hybrid = groq_hybrid_costs(expected_daily_queries)
    _add_table_slide(prs, "Option C: Groq API Hybrid",
        ["Component", "Details", "Monthly Cost"],
        [
            ["Groq API", f"~{expected_daily_queries * 30:,} queries/month", f"${hybrid['breakdown']['Groq API']:,.2f}"],
            ["Cloud VM", "Backend + ChromaDB", f"${hybrid['breakdown']['Cloud VM']:,.0f}"],
            ["Managed PostgreSQL", "Small instance", f"${hybrid['breakdown']['Managed PostgreSQL']:,.0f}"],
            ["Monthly Total", "", f"${hybrid['monthly_cost']:,.2f}"],
            ["Annual Total", "", f"${hybrid['yearly'].year_1:,.2f}"],
            ["5-Year TCO", "", f"${hybrid['yearly'].five_year_tco:,.2f}"],
        ],
    )

    # Slide 14: Implementation Phases
    _add_content_slide(prs, "Implementation Phases", [
        "Phase 1: Backend foundation - FastAPI, PostgreSQL, JWT auth, 5-level access",
        "Phase 2: RAG engine - document processing, embeddings, ChromaDB, retrieval",
        "Phase 3: API endpoints - chat, documents, collections, admin",
        "Phase 4: Frontend - Streamlit UI consuming REST API",
        "Phase 5: Deployment - Docker, SSL, monitoring, backups",
        "Phase 6: SSO integration - connect to university identity provider",
    ])

    # Slide 15: Recommendation
    rec = (
        "Groq API Hybrid" if preferred_deployment == "hybrid"
        else "On-Premise GPU" if preferred_deployment == "onprem"
        else "Cloud (Azure/AWS)" if preferred_deployment == "cloud"
        else "Groq API Hybrid for initial deployment, transitioning to on-premise as usage grows"
    )
    _add_content_slide(prs, "Recommendation", [
        f"Recommended approach: {rec}",
        "",
        "Start with Groq API Hybrid for lowest upfront cost and fastest deployment",
        "Monitor actual usage patterns and query volumes over 3-6 months",
        "Evaluate data sovereignty requirements with university legal team",
        "Plan transition to on-premise if usage exceeds cost-effectiveness threshold",
        "The pluggable LLM provider architecture allows switching without code changes",
    ])

    # Slide 16: Next Steps
    _add_content_slide(prs, "Next Steps", [
        "1. Review and approve deployment option",
        "2. Provision infrastructure (cloud credits or hardware procurement)",
        "3. Configure SSO integration with university identity provider",
        "4. Identify initial document collections and access policies",
        "5. Run pilot with selected department(s)",
        "6. Gather feedback and iterate",
        "7. University-wide rollout",
    ])

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
